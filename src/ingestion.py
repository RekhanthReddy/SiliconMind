"""
SiliconMind Document Ingestion
───────────────────────────────
Handles PDF, TXT, DOCX, PPTX and URL ingestion into ChromaDB.
"""

import io
import re
import urllib.request


def extract_text(file) -> tuple[str, str]:
    """
    Extract text from an uploaded file.
    Returns (text, filename)
    """
    filename = file.name
    content  = ""

    if filename.endswith(".pdf"):
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(file.read())) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        content += text + "\n"
        except ImportError:
            content = file.read().decode("utf-8", errors="ignore")

    elif filename.endswith(".docx"):
        try:
            import docx
            doc = docx.Document(io.BytesIO(file.read()))
            content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            content = file.read().decode("utf-8", errors="ignore")

    elif filename.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file.read()))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += shape.text + "\n"
        except ImportError:
            content = file.read().decode("utf-8", errors="ignore")

    else:
        content = file.read().decode("utf-8", errors="ignore")

    return content.strip(), filename


def extract_from_url(url: str) -> tuple[str, str]:
    """
    Fetch and extract text from a URL (ArXiv abstract page or plain HTML).
    Returns (text, source_label)
    """
    try:
        # ArXiv: use API for clean abstract
        if "arxiv.org/abs/" in url:
            arxiv_id = url.strip("/").split("/")[-1]
            api_url  = f"http://export.arxiv.org/api/query?id_list={arxiv_id}"
            req      = urllib.request.Request(api_url, headers={"User-Agent": "SiliconMind/1.0"})
            with urllib.request.urlopen(req, timeout=8) as resp:
                import xml.etree.ElementTree as ET
                root  = ET.fromstring(resp.read().decode("utf-8"))
                ns    = {"atom": "http://www.w3.org/2005/Atom"}
                entry = root.find("atom:entry", ns)
                if entry:
                    title   = entry.findtext("atom:title",   "", ns).strip()
                    summary = entry.findtext("atom:summary", "", ns).strip()
                    authors = [a.findtext("atom:name", "", ns)
                               for a in entry.findall("atom:author", ns)]
                    text = f"Title: {title}\nAuthors: {', '.join(authors[:5])}\nAbstract: {summary}"
                    return text, f"ArXiv:{arxiv_id}"

        # Generic URL: fetch HTML and strip tags
        req = urllib.request.Request(url, headers={"User-Agent": "SiliconMind/1.0"})
        with urllib.request.urlopen(req, timeout=10) as resp:
            html = resp.read().decode("utf-8", errors="ignore")
        # Strip HTML tags
        text = re.sub(r"<[^>]+>", " ", html)
        text = re.sub(r"\s+", " ", text).strip()
        return text[:5000], url

    except Exception as e:
        return "", f"Error fetching {url}: {e}"


def smart_chunk(text: str, filename: str,
                chunk_size: int = 400, overlap: int = 50) -> list[dict]:
    """
    Split text into overlapping chunks with section-aware metadata.
    Detects headings to add topic tags to each chunk.
    """
    # Detect section headings
    heading_pattern = re.compile(
        r"^(\d+[\.\d]*\s+.{5,80}|[A-Z][A-Z\s]{4,40})$", re.MULTILINE
    )

    # Topic tag detection
    topic_keywords = {
        "ATPG":        ["atpg", "fault coverage", "stuck-at", "transition fault",
                        "pattern generation", "fault model", "test pattern"],
        "Scan Chain":  ["scan chain", "scan cell", "scan flip-flop", "shift register",
                        "scan-in", "scan-out", "edt", "compression"],
        "MBIST":       ["mbist", "memory bist", "march", "bist", "memory test",
                        "sram", "memory repair", "redundancy"],
        "JTAG":        ["jtag", "boundary scan", "tap controller", "tdi", "tdo",
                        "ieee 1149", "ijtag", "ieee 1687"],
        "Silicon Validation": ["bring-up", "silicon", "tester", "ate", "correlation",
                               "validation", "characterisation"],
    }

    words        = text.split()
    chunks       = []
    current_head = ""

    for i in range(0, len(words), chunk_size - overlap):
        chunk_words = words[i:i + chunk_size]
        chunk_text  = " ".join(chunk_words)

        if len(chunk_text.strip()) < 50:
            continue

        # Detect topic
        chunk_lower = chunk_text.lower()
        topics      = [
            t for t, kws in topic_keywords.items()
            if any(kw in chunk_lower for kw in kws)
        ]

        chunks.append({
            "text":     chunk_text,
            "source":   filename,
            "topics":   ", ".join(topics) if topics else "General",
            "chunk_idx": i // (chunk_size - overlap)
        })

    return chunks


def index_to_chromadb(chunks: list[dict], vector_store) -> int:
    """Upsert chunks into ChromaDB. Returns number of chunks stored."""
    if not chunks:
        return 0

    texts     = [c["text"]   for c in chunks]
    ids       = [f"{c['source'].replace(' ','_')}__chunk_{c['chunk_idx']}" for c in chunks]
    metadatas = [{"source": c["source"], "topics": c["topics"]} for c in chunks]

    batch = 50
    for i in range(0, len(texts), batch):
        vector_store.upsert(
            documents=texts[i:i + batch],
            ids=ids[i:i + batch],
            metadatas=metadatas[i:i + batch]
        )
    return len(chunks)
